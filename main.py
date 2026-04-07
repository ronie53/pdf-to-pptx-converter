import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def pdf_to_pptx(pdf_file, output_file):
    if not os.path.exists(pdf_file):
        print(f"Error: File '{pdf_file}' not found.")
        return

    try:
        with pdfplumber.open(pdf_file) as pdf:
            presentation = Presentation()

            for page_num, page in enumerate(pdf.pages):
                print(f"Processing page {page_num + 1}...")
                
                # Get text with character details
                chars = page.chars
                if not chars:
                    continue

                # Sort characters by position
                chars.sort(key=lambda c: (c['top'], c['x0']))

                # Group into lines
                lines = []
                current_line = []
                current_top = None

                for char in chars:
                    if current_top is None or abs(char['top'] - current_top) > 2:
                        if current_line:
                            lines.append(current_line)
                        current_line = [char]
                        current_top = char['top']
                    else:
                        current_line.append(char)

                if current_line:
                    lines.append(current_line)

                # Create slides with detected formatting
                slide = None
                content_frame = None
                line_count = 0
                max_lines = 30

                for line_chars in lines:
                    # Create new slide if needed
                    if line_count >= max_lines:
                        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(6.9))
                        content_frame = content_box.text_frame
                        content_frame.word_wrap = True
                        line_count = 0

                    if slide is None:
                        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(6.9))
                        content_frame = content_box.text_frame
                        content_frame.word_wrap = True

                    # Add line with formatting
                    paragraph = content_frame.add_paragraph()
                    
                    for char in line_chars:
                        run = paragraph.add_run()
                        run.text = char.get('text', '')

                        # Detect bold and font size
                        fontname = char.get('fontname', '').lower()
                        size = char.get('size', 11) or 11
                        is_bold = 'bold' in fontname or size > 12

                        # Use PDF font size as base and boost slightly for readability
                        ppt_size = Pt(size + 2)
                        run.font.size = ppt_size
                        run.font.bold = is_bold

                    paragraph.space_after = Pt(6)
                    line_count += 1

            presentation.save(output_file)
            print(f"✓ Done! Saved as {output_file}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    pdf_file = input("Enter PDF file name: ")
    output_file = pdf_file.replace('.pdf', '.pptx')
    pdf_to_pptx(pdf_file, output_file)
